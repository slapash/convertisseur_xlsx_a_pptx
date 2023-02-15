import os
import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches

# Set the path for the directory containing the Excel files
excel_dir = r'C:\Users\benai\Desktop\xlsx'

# Set the path for the directory to save the PowerPoint files
pptx_dir = r'C:\Users\benai\Desktop\pptx'

# Loop through all the Excel files in the directory
def convert_excel_to_pptx(excel_dir, pptx_dir):
    for excel_file in os.listdir(excel_dir):
        if excel_file.endswith('.xlsx'):
            # Load the Excel file into a pandas dataframe
            df = pd.read_excel(os.path.join(excel_dir, excel_file), None, header= None)
            
            # Loop through each sheet in the Excel file
            for sheet_name in df.keys():
                # Create a new PowerPoint presentation
                prs = Presentation()
                print(sheet_name)
                # Add a slide for each row in the sheet
                for i in range(len(df[sheet_name])):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])

                    # Add the row data to the slide
                    row_data = df[sheet_name].iloc[i]
                    for j, val in enumerate(row_data.dropna()):
                        
                        left = Inches(j * 1.5)
                        top = Inches(1.5)
                        width = Inches(1)
                        height = Inches(1)
                        shape = slide.shapes.add_textbox(
                            
                            left=left,
                            top=top,
                            width=width,
                            height=height
                        )
                        shape.text = str(val)
                
                # Save the PowerPoint file with the same name as the Excel sheet
                pptx_file = os.path.join(pptx_dir, f'{excel_file}_{sheet_name}.pptx')
                prs.save(pptx_file)